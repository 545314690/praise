

import pptx

def read_PPTX(file_path):  #pptx,pptm,
    presentation = pptx.Presentation(file_path)
    results = ''
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        results = results + run.text
    return results.encode('GBK', 'ignore').decode('GBK')

file_path = r'/Users/lisenmiao/Desktop/ppt/zm-修改后2022-10-9/120主我单属你.pps'
print(read_PPTX(file_path))

